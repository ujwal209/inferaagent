import os
import logging
import csv
import PyPDF2
from io import BytesIO, StringIO
import docx
from pptx import Presentation
import cloudinary
import cloudinary.uploader
from supabase import create_client, Client
from langchain_text_splitters import RecursiveCharacterTextSplitter
from dotenv import load_dotenv

from google import genai
from google.genai import types

# 🚀 Load environment variables
load_dotenv()

# 🚀 Cloudinary Setup
cloudinary.config(
    cloud_name=os.getenv("CLOUDINARY_CLOUD_NAME"),
    api_key=os.getenv("CLOUDINARY_API_KEY"),
    api_secret=os.getenv("CLOUDINARY_API_SECRET"),
    secure=True
)

# Safely grab the Supabase keys
SUPABASE_URL = os.getenv("SUPABASE_URL") or os.getenv("NEXT_PUBLIC_SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_KEY") or os.getenv("SUPABASE_SERVICE_ROLE_KEY")

if not SUPABASE_URL or not SUPABASE_KEY:
    logging.error("🚨 CRITICAL: Supabase URL or Key is missing from .env! RAG features are disabled.")
    supabase = None
else:
    supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# 🚀 GEMINI API SETUP (Comma Separated Fallback Keys)
gemini_env = os.getenv("GEMINI_API_KEYS", os.getenv("GEMINI_API_KEY", ""))
GEMINI_API_KEYS = [k.strip() for k in gemini_env.split(",") if k.strip()]


def get_gemini_embedding(text: str, task_type: str = "RETRIEVAL_DOCUMENT") -> list[float]:
    """Uses Google's new SDK to generate 3072-dimensional embeddings with multi-key fallback."""
    if not GEMINI_API_KEYS:
        raise Exception("GEMINI_API_KEYS missing from your .env file!")

    clean_text = text.replace("\n", " ").strip()
    last_error = None

    # Fallback loop: try keys until one works
    for api_key in GEMINI_API_KEYS:
        try:
            client = genai.Client(api_key=api_key)
            result = client.models.embed_content(
                model="gemini-embedding-001",
                contents=clean_text,
                config=types.EmbedContentConfig(task_type=task_type)
            )
            return result.embeddings[0].values
        except Exception as e:
            logging.warning(f"⚠️ Gemini Key ending in ...{api_key[-4:]} failed: {e}")
            last_error = str(e)
            
    raise Exception(f"All Gemini API keys failed! Last error: {last_error}")


def upload_to_cloudinary(file_bytes: bytes, filename: str) -> str:
    """Uploads file to Cloudinary and returns the secure URL."""
    try:
        upload_result = cloudinary.uploader.upload(
            file_bytes,
            resource_type="auto",
            public_id=f"inferacore/{filename}",
            unique_filename=True
        )
        return upload_result["secure_url"]
    except Exception as e:
        logging.error(f"Cloudinary upload failed: {e}")
        raise e


def process_and_store_document(file_bytes: bytes, session_id: str, filename: str) -> str:
    """Uploads to Cloudinary, extracts text, chunks it, and stores it in Supabase using Gemini Embeddings."""
    
    # 1. Upload to Cloudinary (for UI display & Vision Models)
    file_url = upload_to_cloudinary(file_bytes, filename)
    
    full_text = ""
    ext = filename.lower().split('.')[-1]
    
    try:
        # 🚀 STRICT SKIP FOR IMAGES: Do not chunk/embed images! The vision model handles them.
        if ext in ['png', 'jpg', 'jpeg', 'webp', 'gif']:
            logging.info(f"🚀 [RAG] Strictly skipping embedding for image: {ext}. Vision model will handle it.")
            return file_url

        if ext == 'pdf':
            reader = PyPDF2.PdfReader(BytesIO(file_bytes))
            full_text = "".join(page.extract_text() for page in reader.pages)
        elif ext in ['docx', 'doc']:
            doc = docx.Document(BytesIO(file_bytes))
            full_text = "\n".join([para.text for para in doc.paragraphs])
        elif ext in ['pptx', 'ppt']:
            prs = Presentation(BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
        elif ext in ['txt', 'md', 'py', 'js', 'ts', 'json', 'csv']:
            full_text = file_bytes.decode('utf-8')
        else:
            logging.info(f"Skipping RAG embedding for unsupported file type: {ext}")
            return file_url

        if not full_text.strip():
            logging.warning(f"Could not extract any text from {filename}. Skipping RAG.")
            return file_url

        # 3. Break the text into chunks
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
        chunks = splitter.split_text(full_text)

        # 4. Embed and save each chunk to Supabase
        if supabase:
            for chunk in chunks:
                # Use RETRIEVAL_DOCUMENT when writing to DB
                vector = get_gemini_embedding(chunk, task_type="RETRIEVAL_DOCUMENT")
                supabase.table("document_chunks").insert({
                    "session_id": session_id,
                    "content": chunk,
                    "embedding": vector
                }).execute()
        
        return file_url
        
    except Exception as e:
        logging.error(f"Processing error for {filename}: {e}")
        return file_url


def retrieve_relevant_context(query: str, session_id: str) -> str:
    """Searches Supabase for the most relevant document chunks based on the user's query."""
    
    # 1. Skip if no query or no database connection
    if not supabase or not query or not query.strip():
        return ""
        
    try:
        # 🚀 CRITICAL FIX: The Fast Pre-Check!
        # Does this session actually have any uploaded documents?
        has_docs = supabase.table("document_chunks").select("id").eq("session_id", session_id).limit(1).execute()
        
        if not has_docs.data:
            # No documents found for this session! Fast exit.
            return ""

        # 2. ONLY if docs exist, we spend the time/API call to embed the user's question
        query_vector = get_gemini_embedding(query, task_type="RETRIEVAL_QUERY")

        # 3. Perform the vector similarity search
        response = supabase.rpc("match_document_chunks", {
            "query_embedding": query_vector,
            "match_threshold": 0.4, 
            "match_count": 4,       
            "p_session_id": session_id
        }).execute()

        chunks = response.data
        if not chunks:
            return ""

        context = "\n\n----- \n\n".join([f"Excerpt from document:\n{c['content']}" for c in chunks])
        return f"\n\n[SYSTEM: RELEVANT DOCUMENT CONTEXT FOUND]\n{context}\n[END CONTEXT]\n"
    
    except Exception as e:
        logging.warning(f"RAG Retrieval failed: {e}")
        return ""   